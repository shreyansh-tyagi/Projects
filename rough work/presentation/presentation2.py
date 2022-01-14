import os 
from pptx import Presentation
p=Presentation()
l=p.slide_layouts[1]
def create_slide(filename,content):
    content=content.splitlines()  #each lines of string passed in slides will go in and get splited into a list of string
    for s in content:
        s.split("#")
        t,c=s,s
        slide=p.slides.add_slide(l)
        slide.shapes.title.text=t
        slide.placeholders[1].text=c
    p.save(filename)
    os.startfile(filename)


content="""python programming # This ppt is made by using python code 
chatbot gui # A chatbot is an AI-based software designed to 
interact with humans in their natural languages. . A chatbot is arguably one 
of the best applications of natural language processing. A chatbot user interface (UI) 
is a series of graphical and language elements that allow for human-computer interaction.




"""
filename=create_slide("presentation using python week 2.pptx",content)


