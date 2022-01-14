# python-pptx is a Python library for creating and updating PowerPoint (. pptx) files.
# full form of pptx is microsoft powerpoint open XML
from pptx import Presentation
import os


# create ppt presentation to add slides to
a=Presentation()


# now what should be the layout of each slides
# 0= title slides, 1= title and content ,2= subtitle , 3= quotes and so on
slide1_select=a.slide_layouts[0]
 

# adding this slide to presentation
slide1=a.slides.add_slide(slide1_select)

# creating a placeholder for title
title1=slide1.shapes.title

# insert text 
title1.text='PYTHON PROGRAMMING'

# creating subtitle
subtitle1=slide1.placeholders[1]


#insert text
subtitle1.text='Python is an Emotion\n\n\n\nPresented By: Shreyansh Tyagi'


slide2_select=a.slide_layouts[1]
slide2=a.slides.add_slide(slide2_select)

title2=slide2.shapes.title
title2.text=' DESKTOP ASSISTANT'

bar=slide2.shapes
bar1=bar.placeholders[1]
bar2=bar1.text_frame.add_paragraph()
bar2.text="Desktop Voice Assistant with Voice Recognition Intelligence, which takes the user \
input in form of user's voice and processes it and return the output in various ways like an action to \
be performed or the search result is speaked out to the end user."
bar2.level=3

slide3_select=a.slide_layouts[1]
slide3=a.slides.add_slide(slide3_select)

title3=slide3.shapes.title
title3.text='PACKAGES REQUIREMENT'

bar=slide3.shapes
bar1=bar.placeholders[1]
bar2=bar1.text_frame.add_paragraph()
bar2.text="pip install SpeechRecognition\n \
pip install pyttsx3\n \
pip install playsound\n \
pip install Pillow\n \
pip install pyscreenshot\n \
pip install pynput\n \
pip install psutil \
pip install opencv-contrib-python\n \
pip install opencv\n \
pip install wikipedia \
pip install webbrowser [Available with installer] pip install bs4\n \
pip install smtplib [Available with installer]\n "
bar2.level=3


slide4_select=a.slide_layouts[1]
slide4=a.slides.add_slide(slide3_select)

title4=slide4.shapes.title
title4.text='CHATBOT GUI'


bullet=slide4.shapes
c=bullet.placeholders[1]
c.text='tokenization\nstemming\nbag_of_words'


slide5_select=a.slide_layouts[1]
slide5=a.slides.add_slide(slide5_select)

title2=slide5.shapes.title
title2.text='TOKENIZATION'

bar=slide5.shapes
bar1=bar.placeholders[1]
bar2=bar1.text_frame.add_paragraph()
bar2.text="Tokenization is the process by which a large quantity of text is divided into  \
smaller parts called tokens. These tokens are very useful for finding patterns and are considered  \
as a base step for stemming and lemmatization.\n \
for example:-sentence_data = The First sentence is about Python. The Second: about Django. You can learn Python,Django and Data Ananlysis here.\n\n \
['The First sentence is about Python.', 'The Second: about Django.', 'You can learn Python,Django and Data Ananlysis here.']"
bar2.level=2


slide6_select=a.slide_layouts[1]
slide6=a.slides.add_slide(slide6_select)

title2=slide6.shapes.title
title2.text='STEMMING'

bar=slide6.shapes
bar1=bar.placeholders[1]
bar2=bar1.text_frame.add_paragraph()
bar2.text='Stemming is the process of producing morphological variants of a root/base word. \
Stemming programs are commonly referred to as stemming algorithms or stemmers.\n for example:- words = ["program", "programs", "programmer", "programing", "programmers"]\n\n \
   [program,program,program,program,program]'
bar2.level=2


slide7_select=a.slide_layouts[1]
slide7=a.slides.add_slide(slide7_select)

title2=slide7.shapes.title
title2.text='BAG_OF_WORD'

bar=slide7.shapes
bar1=bar.placeholders[1]
bar2=bar1.text_frame.add_paragraph()
bar2.text=' a Natural Language Processing technique of text modeling known as Bag of Words model. \
 Whenever we apply any algorithm in NLP, it works on numbers. We cannot directly feed our text into \
  that algorithm. Hence, Bag of Words model is used to preprocess the text by converting it into a \
  bag of words, which keeps a count of the total occurrences of most frequently used words.'
bar2.level=2 

slide8_select=a.slide_layouts[1]
slide8=a.slides.add_slide(slide8_select)

title2=slide8.shapes.title
title2.text='WEEK OVERVIEW'

bar=slide8.shapes
bar1=bar.placeholders[1]
bar1.text='Problem solving on leetcode\nDaily contribution on github\nUpdated linkedin profile\nProblem solving on hackerrank'

slide9_select=a.slide_layouts[6]
slide9=a.slides.add_slide(slide9_select)

title2=slide9.shapes
title2.text='THANK YOU'

# saving the file in pptx format and giving it a name
a.save('presentation using python.pptx')
os.startfile('presentation using python.pptx')
